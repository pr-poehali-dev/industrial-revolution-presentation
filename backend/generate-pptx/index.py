import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import io
import base64

def handler(event, context):
    '''
    Business: Generate PPTX presentation about Industrial Revolution
    Args: event with httpMethod
    Returns: PPTX file as base64 or download link
    '''
    method = event.get('httpMethod', 'GET')
    
    if method == 'OPTIONS':
        return {
            'statusCode': 200,
            'headers': {
                'Access-Control-Allow-Origin': '*',
                'Access-Control-Allow-Methods': 'GET, POST, OPTIONS',
                'Access-Control-Allow-Headers': 'Content-Type, X-User-Id',
                'Access-Control-Max-Age': '86400'
            },
            'body': ''
        }
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    industrial_dark = RGBColor(28, 25, 23)
    industrial_copper = RGBColor(180, 83, 9)
    industrial_light = RGBColor(214, 211, 209)
    white = RGBColor(255, 255, 255)
    
    def add_title_slide(title, subtitle):
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = industrial_dark
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = white
        title_para.font.name = 'Oswald'
        
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(0.8))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.alignment = PP_ALIGN.CENTER
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = industrial_light
        subtitle_para.font.name = 'Open Sans'
        
        return slide
    
    def add_content_slide(title, items):
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = industrial_dark
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = industrial_copper
        title_para.font.name = 'Oswald'
        
        y_pos = 1.5
        for item in items:
            card_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(9), Inches(1.6))
            card_frame = card_box.text_frame
            card_frame.word_wrap = True
            
            heading = card_frame.paragraphs[0]
            heading.text = item['heading']
            heading.font.size = Pt(20)
            heading.font.bold = True
            heading.font.color.rgb = industrial_copper
            heading.font.name = 'Oswald'
            
            subheading = card_frame.add_paragraph()
            subheading.text = item['subheading']
            subheading.font.size = Pt(14)
            subheading.font.color.rgb = industrial_light
            subheading.font.name = 'Open Sans'
            
            desc = card_frame.add_paragraph()
            desc.text = item['description']
            desc.font.size = Pt(12)
            desc.font.color.rgb = white
            desc.font.name = 'Open Sans'
            desc.space_before = Pt(6)
            
            y_pos += 1.8
        
        return slide
    
    add_title_slide('ИНДУСТРИАЛЬНАЯ РЕВОЛЮЦИЯ', 'Экономическая трансформация человечества')
    
    stages_items = [
        {
            'heading': 'Первая революция (1760-1840)',
            'subheading': 'Великобритания',
            'description': 'Механизация производства, паровой двигатель, текстильная промышленность. Паровые машины Уатта, железные дороги.'
        },
        {
            'heading': 'Вторая революция (1870-1914)',
            'subheading': 'США, Германия',
            'description': 'Электричество, нефть, массовое производство. Конвейер Форда, электрификация городов, химическая промышленность.'
        },
        {
            'heading': 'Третья революция (1950-2000)',
            'subheading': 'Глобальная',
            'description': 'Компьютеризация, автоматизация, информационные технологии. Роботизация производства, интернет, глобализация экономики.'
        }
    ]
    add_content_slide('Этапы Индустриальной Революции', stages_items)
    
    inventions_items = [
        {
            'heading': 'Паровой двигатель (1769)',
            'subheading': 'Джеймс Уатт',
            'description': 'Революция в транспорте и производстве. Паровозы, пароходы, фабричные станки. Независимый источник энергии.'
        },
        {
            'heading': 'Ткацкий станок (1785)',
            'subheading': 'Эдмунд Картрайт',
            'description': 'Рост производительности в 40 раз. Текстильные фабрики, массовое производство ткани, снижение цен на одежду.'
        },
        {
            'heading': 'Электрическая лампа (1879)',
            'subheading': 'Томас Эдисон',
            'description': 'Круглосуточная работа предприятий. Ночные смены, рост производства вдвое, урбанизация и электрификация городов.'
        }
    ]
    add_content_slide('Ключевые Изобретения', inventions_items)
    
    results_items = [
        {
            'heading': 'Рост ВВП: +400%',
            'subheading': '1800-1900',
            'description': 'Великобритания: с £350M до £2B. США: рост экономики в 15 раз. Германия: промышленный бум и индустриализация.'
        },
        {
            'heading': 'Урбанизация: 10% → 80%',
            'subheading': '1800-2000',
            'description': 'Лондон: 1M → 7M жителей. Новые промышленные города: Манчестер, Бирмингем, Детройт. Развитие инфраструктуры.'
        },
        {
            'heading': 'Производительность труда: +1500%',
            'subheading': '1760-1900',
            'description': 'Текстиль: с 1 до 40 метров/день. Сталь: рост производства в 200 раз. Снижение стоимости товаров и рост уровня жизни.'
        }
    ]
    add_content_slide('Экономические Результаты', results_items)
    
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    
    pptx_base64 = base64.b64encode(output.read()).decode('utf-8')
    
    return {
        'statusCode': 200,
        'headers': {
            'Content-Type': 'application/json',
            'Access-Control-Allow-Origin': '*'
        },
        'isBase64Encoded': False,
        'body': json.dumps({
            'success': True,
            'filename': 'industrial_revolution.pptx',
            'fileBase64': pptx_base64,
            'size': len(pptx_base64)
        })
    }
